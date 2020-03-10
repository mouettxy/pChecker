import json
import re
import string
import random

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.enum.text import MSO_AUTO_SIZE

from PIL import Image, ImageDraw

from globals import TEXT_THRESHOLD


class TextWrapper(object):
    """
    Igor Pomaranskiy
    https://stackoverflow.com/users/535884/igor-pomaranskiy
    Helper class to wrap text in lines, based on given text, font
    and max allowed line width.
    """

    def __init__(self, text, font, max_width):
        self.text = text
        self.text_lines = [
            ' '.join([w.strip() for w in l.split(' ') if w])
            for l in text.split('\n')
            if l
        ]
        self.font = font
        self.max_width = max_width

        self.draw = ImageDraw.Draw(
            Image.new(
                mode='RGB',
                size=(100, 100)
            )
        )

        self.space_width = self.draw.textsize(
            text=' ',
            font=self.font
        )[0]

    def get_text_width(self, text):
        return self.draw.textsize(
            text=text,
            font=self.font
        )[0]

    def wrapped_text(self):
        wrapped_lines = []
        buf = []
        buf_width = 0
        for line in self.text_lines:
            for word in line.split(' '):
                word_width = self.get_text_width(word)
                expected_width = word_width if not buf else \
                    buf_width + self.space_width + word_width
                if expected_width <= self.max_width:
                    # word fits in line
                    buf_width = expected_width
                    buf.append(word)
                else:
                    # word doesn't fit in line
                    wrapped_lines.append(' '.join(buf))
                    buf = [word]
                    buf_width = word_width
            if buf:
                wrapped_lines.append(' '.join(buf))
                buf = []
                buf_width = 0
        return '\n'.join(wrapped_lines)

    def total_width_height(self):
        lines = self.wrapped_text().splitlines()
        width = []
        height = []
        for line in lines:
            width.append(self.font.getsize(line)[0])
            height.append(self.font.getsize(line)[1])
        return max(width), sum(height) + 5


class Utils:
    """
    Класс от которого наследуются все остальные. Его задача получить презентацию, и служить контейнером для различных
    полезных функций.
    """
    def __init__(self, path_to_presentation):
        """
        :param path_to_presentation: abspath к файлу презентации для проверки.
        """
        super().__init__()
        self.presentation = Presentation(path_to_presentation)
        self.path_to_presentation = path_to_presentation

    @staticmethod
    def to_json(data):
        """
        Служит для переноса информации в JSON формат
        :param data: Любая информация
        :type data: dict, list, tuple, str, int, long, float, True, False, None
        :return: JSON формат
        :rtype: object, array, string, number, true, false, null
        """
        return json.dumps(data)

    @staticmethod
    def is_image(shape):
        """
        :param shape: Обьект презентации
        :return: Является ли обьект картинкой
        :rtype: bool
        """
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.PICTURE:
            return True
        return False

    @staticmethod
    def is_title(shape):
        """
        :param shape: Обьект презентации
        :return: Является ли обьект заголовком
        :rtype: bool
        """
        if shape.is_placeholder and (
                shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.TITLE
                or shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.SUBTITLE
                or shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.VERTICAL_TITLE
                or shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.CENTER_TITLE):
            return True
        return False

    def is_text(self, shape):
        """
        :param shape: Обьект презентации
        :return: Является ли обьект текстом
        :rtype: bool
        """
        if hasattr(shape, "text"):
            return True
        if shape.has_text_frame:
            if self.is_title(shape):
                return True
            if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.BODY:
                return True
            if len(shape.text) > TEXT_THRESHOLD:
                return True
        return False

    @staticmethod
    def is_not_auto_size(shape):
        if shape.text_frame.auto_size is None or \
                shape.text_frame.auto_size == MSO_AUTO_SIZE.NONE or \
                shape.text_frame.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT:
            return True
        return False

    @staticmethod
    def string_optimize(data):
        """
        :param data: Строчка/Параграф из презентации
        :type data: str
        :return: Очищенная строка без лишних символов и слов менее 4х букв
        :rtype: str
        """
        delete_junk_symbols = re.compile('[^a-zA-Zа-яА-ЯёЁ]')
        delete_junk_words = re.compile('\\b\\w{0,3}\\b')
        return delete_junk_words.sub("", delete_junk_symbols.sub("", data.lower()))

    @staticmethod
    def convert_emu_px(emu):
        """
        :param emu: Значение в EMU
        :type emu: int
        :return: Округлённое конвертированное значение в пикселях
        :rtype: int
        """
        return round(emu // 9525)

    @staticmethod
    def _translate_results(results):
        """
        :param results: Принимает результат работы analyze_results()
        :return: Dict где ключ переведён из ID в Str
        """
        results['Количество слайдов'] = results.pop(0)
        results['Блоки текста и изображений размещены'] = results.pop(1)
        results['Название на титульном'] = results.pop(2)
        results['Название на 2м и 3м слайде'] = results.pop(3)
        results['Единый шрифт'] = results.pop(4)
        results['Правильный размер шрифта'] = results.pop(5)
        results['Изображения не искажены'] = results.pop(6)
        return results

    @staticmethod
    def check_collision(first, second):
        """
        Проверяет 2 прямоугольника на наличие коллизий между ними.
        :param first: Список состоящий из x, y, width, height
        :type first: list of int or float
        :param second: Список состоящий из x, y, width, height
        :type second: list of int or float
        :return: True если найдена коллизия False иначе
        :rtype: boolean
        """
        if first[0] < second[0] + second[2] and first[0] + first[2] > second[0] and \
                first[1] < second[1] + second[3] and first[1] + first[3] > second[1]:
            return True
        return False

    @staticmethod
    def random_string(symbols=5):
        return ''.join(random.choices(string.ascii_uppercase + string.digits, k=symbols))

