# -*- coding: utf-8 -*-
import base64
import json
import os
import re
from collections import Counter
from pathlib import Path

# import time
import string
import random
import win32com.client
import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from pandas.errors import EmptyDataError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.enum.text import MSO_AUTO_SIZE

from globals import *

# from pp_classes import MSOPPT, MSO

"""
Issues:
https://github.com/scanny/python-pptx/issues/147
https://github.com/scanny/python-pptx/issues/337
"""


class GlobalUtils(object):
    """
    Даёт некоторые полезные функции/фичи для разработки. Возможно будет убран позднее.
    """

    def __init__(self, dir="D:\\Presentations"):
        self.get_list_of_presentations = [dir + "\\" + d for d in os.listdir(dir)]


class TextWrapper(object):
    """
    Igor Pomaranskiy
    https://stackoverflow.com/users/535884/igor-pomaranskiy
    Helper class to wrap text in lines, based on given text, font
    and max allowed line width.
    """

    def __init__(self, text, font, max_width):
        self.text = text
        self.text_lines = [
            ' '.join([w.strip() for w in l.split(' ') if w])
            for l in text.split('\n')
            if l
        ]
        self.font = font
        self.max_width = max_width

        self.draw = ImageDraw.Draw(
            Image.new(
                mode='RGB',
                size=(100, 100)
            )
        )

        self.space_width = self.draw.textsize(
            text=' ',
            font=self.font
        )[0]

    def get_text_width(self, text):
        return self.draw.textsize(
            text=text,
            font=self.font
        )[0]

    def wrapped_text(self):
        wrapped_lines = []
        buf = []
        buf_width = 0
        for line in self.text_lines:
            for word in line.split(' '):
                word_width = self.get_text_width(word)
                expected_width = word_width if not buf else \
                    buf_width + self.space_width + word_width
                if expected_width <= self.max_width:
                    # word fits in line
                    buf_width = expected_width
                    buf.append(word)
                else:
                    # word doesn't fit in line
                    wrapped_lines.append(' '.join(buf))
                    buf = [word]
                    buf_width = word_width
            if buf:
                wrapped_lines.append(' '.join(buf))
                buf = []
                buf_width = 0
        return '\n'.join(wrapped_lines)

    def total_width_height(self):
        lines = self.wrapped_text().splitlines()
        width = []
        height = []
        for line in lines:
            width.append(self.font.getsize(line)[0])
            height.append(self.font.getsize(line)[1])
        return max(width), sum(height) + 5


class PrintTo:
    """
    Класс PrintTo
    Конструктор принимает:
    :param results: Переведённый результат работы  CheckPresentationAnalyze.analyze_rezults()
    :param path_to_output: Полный путь к файлу куда будет выгружен результат
    :param path_to_pptx: Полный путь к файлу презентации
    :param encoding: Кодировка которая должна получиться в выходном файле
    TODO: Добавить поддержку изменения кодировки к PrintTo.txt
    TODO: Добавить генерацию файла Excel PrintTo.excel
    """

    def __init__(self, results, path_to_output, path_to_pptx, encoding='utf-8'):
        self.results = results
        self.path_to_output = Path(path_to_output)
        self.path_to_pptx = Path(path_to_pptx)
        self.encoding = encoding
        self.output_name, self.output_extension = os.path.splitext(self.path_to_output)
        self.results_keys = []
        self.results_values = []
        for result in self.results:
            self.results_keys.append(result)
            self.results_values.append(self.results[result])
        self.results_zip = list(zip(self.results_keys, self.results_values))
        self.mode_list = ['write', 'rewrite']

    def _extension_check(self, extension):
        """
        :param extension: Расширение файла str (пример: ".txt")
        :return: None в случае когда расширение совпадает с вызванное функцией, иначе генерирует подробную ошибку
        """
        if self.output_extension != extension:
            raise Exception('Ошибка! Неверное расширение файла. \n'
                            f'Ожидалось "{extension}", введено "{self.output_extension}". \n'
                            'Проверьте правильность введённого пути к конечному файлу.')
        return

    def _write_mode_check(self, mode):
        """
        :param mode: Режим записи в файл из значений self.mode_list
        :return: Возвращает str применимую к режиму записи файла
        """
        mode_string = ', '.join(self.mode_list)
        if mode == 'write':
            return 'a'
        elif mode == 'rewrite':
            return 'w'
        else:
            raise Exception('Ошибка! Неверно указан метод открытия файла.\n'
                            f'Ожидаемые значения "{mode_string}", введено "{mode}". \n'
                            'Доступные методы: \n'
                            ' write - файл не будет перезаписан, полученный результат добавиться к предыдущему. \n'
                            ' rewrite - файл будет перезаписан, данные которые были до этого будут удалены. \n')

    def _empty(self, file_path, extension):
        """
        :param file_path: Путь к файлу (пример: 'C:/path/to/file.extension')
        :param extension: Расширение файла (пример: '.txt')
        :return: True если файл пуст, False если файл не пуст
        """
        if extension == '.txt':
            if os.stat(file_path).st_size > 0:
                return False
            return True
        elif extension == '.csv':
            try:
                file_contents = pd.read_csv(file_path)
                return file_contents.empty
            except EmptyDataError:
                return True
            except (OSError, IOError):
                create_file = open(file_path, 'w')
                create_file.close()
                return True
        elif extension == '.xlsx':
            try:
                file_contents = pd.read_excel(file_path)
                return file_contents.empty
            except EmptyDataError:
                return True
            except (OSError, IOError):
                create_file = open(file_path, 'w')
                create_file.close()
                return True

    def _write_to_csv(self, data, file_path, mode, encoding):
        """
        :param data: Список для заполнения
        :param file_path: Путьк файлу
        :param mode: Режим записи в файл
        :param encoding: Кодировка текста внутри
        :return: None или Str в случае выполнения, игаче генерирует подробные ошибки.
        """
        data_frame = pd.DataFrame(data)
        try:
            return data_frame.to_csv(file_path, mode=mode, header=True, encoding=encoding, index=False)
        except PermissionError:
            raise Exception('Недостаточно прав для открытия, создания, или записи в файл. \n'
                            'Попробуйте закрыть файл, или проверить права на запись и чтение файла. \n')

    def txt(self, mode):
        """
        :param mode: Режим записи в файл
        :return: Возвращает str в случае успешного выполнения, иначе генерирует подробные ошибки
        """
        self._extension_check('.txt')
        mode = self._write_mode_check(mode)
        txt_file = open(self.path_to_output, mode=mode, encoding='utf-8')
        if not (self._empty(self.path_to_output, '.txt')):
            txt_file.write('\n')
        txt_file.write(f'Проверка файла: {self.path_to_pptx}\n')
        for result in self.results_zip:
            if result[0] == 'Количество слайдов':
                txt_file.write(f'{result[0]} => {result[1]}\n')
            else:
                txt_file.write(f'{result[0]} => {"Да" if result[1] else "Нет"}\n')
        txt_file.close()
        return 'Успешная запись в файл'

    def csv(self, mode):
        """
        :param mode: Режим записи в файл
        :return: Возвращает str в случае успешного выполнения, иначе генерирует подробные ошибки
        """
        self._extension_check('.csv')
        mode = self._write_mode_check(mode)
        data_without_columns = [self.results_values]
        data_with_columns = [self.results]
        if self._empty(self.path_to_output, '.csv') or (not (self._empty(self.path_to_output, '.csv')) and mode == 'w'):
            self._write_to_csv(data_with_columns, self.path_to_output, mode, self.encoding)
            return 'Успешная запись в файл'
        elif not (self._empty(self.path_to_output, '.csv')) and mode == 'a':
            self._write_to_csv(data_without_columns, self.path_to_output, mode, self.encoding)
            return 'Успешная запись в файл'

    def excel(self, mode):
        pass


class CheckPresentation:
    def __init__(self, path):
        super().__init__()
        self.presentation = Presentation(path)
        self.path_to_presentation = path
        self.text_threshold = TEXT_THRESHOLD


class CheckPresentationUtils(CheckPresentation):
    """
    Различные полезные методы. От этого гласса наследуются CheckPresentationGetData и CheckPresentationAnalyze
    """

    @staticmethod
    def _to_json(string):
        """
        Служит для переноса информации в JSON формат
        :param string: Любая информация
        :type string: dict, list, tuple, str, int, long, float, True, False, None
        :return: JSON формат
        :rtype: object, array, string, number, true, false, null
        """
        return json.dumps(string)

    @staticmethod
    def is_image(shape):
        """
        :param shape: Обьект презентации
        :return: Является ли обьект картинкой
        :rtype: bool
        """
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.PICTURE:
            return True
        return False

    @staticmethod
    def is_title(shape):
        """
        :param shape: Обьект презентации
        :return: Является ли обьект заголовком
        :rtype: bool
        """
        if shape.is_placeholder and (
                shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.TITLE
                or shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.SUBTITLE
                or shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.VERTICAL_TITLE
                or shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.CENTER_TITLE):
            return True
        return False

    def is_text(self, shape):
        """
        :param shape: Обьект презентации
        :return: Является ли обьект текстом
        :rtype: bool
        """
        if hasattr(shape, "text"):
            return True
        if shape.has_text_frame:
            if self.is_title(shape):
                return True
            if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.BODY:
                return True
            if len(shape.text) > self.text_threshold:
                return True
        return False

    @staticmethod
    def is_not_auto_size(shape):
        if shape.text_frame.auto_size is None or \
                shape.text_frame.auto_size == MSO_AUTO_SIZE.NONE or \
                shape.text_frame.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT:
            return True
        return False

    @staticmethod
    def string_optimize(string):
        """
        :param string: Строчка/Параграф из презентации
        :type string: str
        :return: Очищенная строка без лишних символов и слов менее 4х букв
        :rtype: str
        """
        delete_junk_symbols = re.compile('[^a-zA-Zа-яА-ЯёЁ]')
        delete_junk_words = re.compile('\\b\\w{0,3}\\b')
        return delete_junk_words.sub("", delete_junk_symbols.sub("", string.lower()))

    @staticmethod
    def convert_emu_px(emu):
        """
        :param emu: Значение в EMU
        :type emu: int
        :return: Округлённое конвертированное значение в пикселях
        :rtype: int
        """
        return round(emu // 9525)

    @staticmethod
    def _translate_results(results):
        """
        :param results: Принимает результат работы analyze_results()
        :return: Dict где ключ переведён из ID в Str
        """
        results['Количество слайдов'] = results.pop(0)
        results['Блоки текста и изображений размещены'] = results.pop(1)
        results['Название на титульном'] = results.pop(2)
        results['Название на 2м и 3м слайде'] = results.pop(3)
        results['Единый шрифт'] = results.pop(4)
        results['Правильный размер шрифта'] = results.pop(5)
        results['Изображения не искажены'] = results.pop(6)
        return results

    @staticmethod
    def check_collision(first, second):
        """
        Проверяет 2 прямоугольника на наличие коллизий между ними.
        :param first: Список состоящий из x, y, width, height
        :type first: list of int or float
        :param second: Список состоящий из x, y, width, height
        :type second: list of int or float
        :return: True если найдена коллизия False иначе
        :rtype: boolean
        """
        if first[0] < second[0] + second[2] and first[0] + first[2] > second[0] and \
                first[1] < second[1] + second[3] and first[1] + first[3] > second[1]:
            return True
        return False

    @staticmethod
    def random_string(symbols=5):
        return ''.join(random.choices(string.ascii_uppercase + string.digits, k=symbols))


class CheckPresentationGetData(CheckPresentationUtils):
    """
    Получает различные данные из презентации
    """

    def length(self):
        """
        :return: Количество слайдов
        :rtype: int
        """
        return len(self.presentation.slides)

    def slide_by_id(self, slide_id):
        """
        :param slide_id: ID слайда
        :return: Слайд
        :rtype: class
        """
        return self.presentation.slides.get(slide_id)

    def prs_w_h(self):
        """
        :return: Значение ширины и длины слайдов презентации в пикселях.
        :rtype: tuple of (int, int)
        """
        return self.convert_emu_px(self.presentation.slide_width), self.convert_emu_px(self.presentation.slide_height)

    def font_sizes_by_shape(self, shape, flag):
        """
        Возвращает размер шрифта в shape, если размер шрифта не найден, то получаем размер шрифта по умолчанию.
        ================================================================================================================
        :param shape: shape obj/class презентации
        :param flag: true/false использовать/не использовать проверку наличия размера шрифта
        :return font_sizes: Размеры шрифта в shape
        :rtype tuple of float
        """
        font_sizes = []
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                try:
                    if self.is_not_auto_size(shape):
                        font_sizes.append(run.font.size.pt)
                    else:
                        font_scale = shape.text_frame._bodyPr.normAutofit.fontScale
                        font_sizes.append(run.font.size.pt * font_scale / 100)
                except AttributeError:
                    pass
        if flag:
            if len(font_sizes) == 0:
                font_sizes.append(18.0)
        return font_sizes

    def text_blocks(self):
        """
        :return: dict of list of (tuple of (int, int), tuple of (int, int), string)
        """
        text_on_slides = {
            1: [],
            2: [],
            3: [],
        }
        for slide in self.presentation.slides:
            slide_index = int(self.presentation.slides.index(slide) + 1)
            for shape in slide.shapes:
                if self.is_text(shape):
                    font_size = self.font_sizes_by_shape(shape, True)
                    left_top = (self.convert_emu_px(shape.left), self.convert_emu_px(shape.top))
                    width_height = (self.convert_emu_px(shape.width), self.convert_emu_px(shape.height))
                    shape_text = shape.text_frame.text.strip()
                    text_on_slides[slide_index].append([left_top, width_height, shape_text, font_size])
        return text_on_slides

    def shapes_by_slide_id(self, slide_id):
        """
        :param slide_id: ID слайда
        :return shapes: Обьекты презентации
        :rtype: list of class
        """
        shapes = []
        for shape in self.slide_by_id(slide_id).shapes:
            shapes.append(shape)
        return shapes

    def font_sizes_by_slide_id(self, slide_id):
        """
        :param slide_id: ID слайда
        :return font_sizes: Все существубщие размеры текста в слайде
        :rtype font_sizes: tuple of float
        """
        font_sizes = []
        for shape in self.shapes_by_slide_id(slide_id):
            if self.is_text(shape):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        try:
                            if self.is_not_auto_size(shape):
                                font_sizes.append(run.font.size.pt)
                            else:
                                font_scale = shape.text_frame._bodyPr.normAutofit.fontScale
                                font_sizes.append(run.font.size.pt * font_scale / 100)
                        except AttributeError:
                            pass
        if len(font_sizes) == 0:
            font_sizes.append(18.0)
        return font_sizes

    def paragraph_runs(self):
        """
        :return runs: "Runs" Параграфов в презентации
        :rtype runs: list of class
        """
        runs = []
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if self.is_text(shape):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            runs.append(run)
        return runs

    def get_font_sizes(self):
        """
        :return font_sizes: Размеры текста в каждом слайде
        :rtype font_sizes: dict of (int, list of float)
        """
        font_sizes = {}
        for slide in self.presentation.slides:
            slide_index = int(self.presentation.slides.index(slide) + 1)
            font_sizes.update({slide_index: self.font_sizes_by_slide_id(slide.slide_id)})
        return font_sizes

    @staticmethod
    def get_typefaces_by_shape(shape):
        """
        Получает названия шрифта по shape презентации
        :param shape: shape презентации
        :type shape: class
        :return: Строку с названием шрифта
        :rtype: string
        """
        typeface = ""
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                try:
                    typeface = run.font.name
                except AttributeError:
                    pass
        if typeface is None:
            typeface = DEFAULT_FONT
        if typeface == "+mn-lt":
            typeface = DEFAULT_THEME_FONT
        return typeface

    def get_typefaces(self):
        """
        :return typefaces: Уникальные название шрифта из презентации
        :rtype typefaces: set
        """
        typefaces = set()
        for run in self.paragraph_runs():
            try:
                typefaces.add(run.font.name)
            except AttributeError:
                pass
        typefaces = list(typefaces)
        for i in range(len(typefaces)):  # На основе темы презентации по умолчанию
            if typefaces[i] is None and DEFAULT_FONT not in typefaces:
                typefaces[i] = DEFAULT_FONT
            else:
                typefaces[i] = DEFAULT_THEME_FONT
            if typefaces[i] == "+mn-lt":
                typefaces[i] = DEFAULT_THEME_FONT
        return typefaces

    def get_text(self):
        """
        :return text: Весь текст в презентации
        :rtype text: list of str
        """
        text = []
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if self.is_text(shape):
                    text.append(shape.text)
        return text

    def analyze_text(self):
        """
        :return most_common: Топ 5 слов в презентации
        :rtype most_common: list of tuple of (str, int)
        """
        text_analyzed = []
        for text in self.get_text():
            text_analyzed.extend(self.string_optimize(text).split())
        most_common = Counter(text_analyzed).most_common(5)
        return most_common

    def dimensions_to_draw(self, shape):
        """
        Возвращает словарь измерений для обьекта презентации.
        :param shape: shape презентации
        :type shape: class
        :return: Словарь изменений для обьекта презентаций содержащий ключи x1, y1, x2, y2, width, height
        :rtype: dict string/int
        """
        width, height = self.convert_emu_px(shape.width), self.convert_emu_px(shape.height)
        left, top = self.convert_emu_px(shape.left), self.convert_emu_px(shape.top)
        x1, y1 = left, top
        x2, y2 = width + x1, height + y1
        return {
            'x1': int(x1),
            'x2': int(x2),
            'y1': int(y1),
            'y2': int(y2),
            'width': int(width),
            'left': int(left),
            'top': int(top),
            'height': int(height),
        }


class CheckPresentationAnalyze(CheckPresentationGetData):
    def image_overlaps(self):
        """
        Проверяет обьекты презентации на коллизии.
        ================================================================================================================
        WIP
        ================================================================================================================
        :return: Список с строками вида "Коллизия между: обьект1 и обьект2 на номер_слайда слайде
        :rtype: list of string
        """
        result = []
        for slide in self.presentation.slides:
            index = int(self.presentation.slides.index(slide) + 1)
            previous_dimensions = []
            previous_shape = None
            for shape in slide.shapes:
                dim = self.dimensions_to_draw(shape)
                current_dimensions = [dim['x1'], dim['y1'], dim['width'], dim['height']]
                if len(previous_dimensions) > 0 and previous_shape is not None:
                    if self.check_collision(current_dimensions, previous_dimensions):
                        result.append(f"Коллизия между: {shape.name} и {previous_shape.name} на {index} слайде")
                previous_dimensions = current_dimensions
                previous_shape = shape
        return result

    def distorted_images(self):
        """
        Проверяет пропорциональность изображений учитывая возожную погрешность до 20 пикселей в каждую сторону.
        :return: True если хоть одна картинка не пропорциональна, False иначе
        :return: boolean
        """
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if self.is_image(shape):
                    tmp = open('tmp.jpg', 'wb')
                    tmp.write(base64.b64decode(base64.b64encode(shape.image.blob)))
                    tmp.close()
                    tmp_pil = Image.open('tmp.jpg')
                    # ширина и высота картинки в презентации
                    p_w, p_h = self.convert_emu_px(shape.width), self.convert_emu_px(shape.height)
                    # ширина и высота оригинальной картинки
                    o_w, o_h = tmp_pil.size
                    tmp_pil.close()
                    os.remove('tmp.jpg')
                    need_width, need_height = o_w * (p_h / o_h), o_h * (p_w / o_w)
                    diff_width, diff_height = abs(p_w - need_width), abs(p_h - need_height)
                    if diff_width > POSSIBLE_DISTORTED_IMAGE_VALUE or diff_height > POSSIBLE_DISTORTED_IMAGE_VALUE:
                        return True
        return False

    def get_slides_contents(self):
        """
        :return slides: Возвращает количество блоков текста, изображений, названий, в каждом слайде с 1-3 включительно.
        """
        slides = {
            1: {
                'textCounter': 0,
                'pictureCounter': 0,
                'titleCounter': 0,
            },
            2: {
                'textCounter': 0,
                'pictureCounter': 0,
                'titleCounter': 0,
            },
            3: {
                'textCounter': 0,
                'pictureCounter': 0,
                'titleCounter': 0,
            }
        }
        for slide in self.presentation.slides:
            index = int(self.presentation.slides.index(slide) + 1)
            for shape in slide.shapes:
                if self.is_text(shape):
                    if self.is_title(shape):
                        slides[index]['titleCounter'] += 1
                    else:
                        slides[index]['textCounter'] += 1
                if self.is_image(shape):
                    slides[index]['pictureCounter'] += 1
        return slides

    def analyze_results(self):
        """
        :return: Dict где ключ => ID критерия, значение => Int or Boolean
        .. note:: ID критериев:
            Int 0 => Str 'Количество слайдов'
            Int 1 => Str 'Блоки текста и изображений размещены'
            Int 2 => Str 'Название на титульном'
            Int 3 => Str 'Название на 2м и 3м слайде'
            Int 4 => Str 'Единый шрифт'
            Int 5 => Str 'Правильный размер шрифта'
            Int 6 => Str 'Изображения не искажены'
        """
        slides_count = self.length()
        if slides_count < 3 or slides_count > 4:
            error = {
                'Ошибка': 'Количество слайдов менее или более трёх.',
            }
            return error

        analyze_params = {
            0: slides_count,
            1: None,
            2: None,
            3: None,
            4: None,
            5: None,
            6: self.distorted_images(),
        }
        slides_contents = self.get_slides_contents()
        font_sizes = self.get_font_sizes()
        typefaces = self.get_typefaces()
        # Check first slide #
        slide1_font_size = 0
        slide1_blocks_correct = 0
        if max(font_sizes[1]) == 40 and min(font_sizes[1]) == 24:
            slide1_font_size = 1
        if slides_contents[1]['titleCounter'] + slides_contents[1]['textCounter'] == 2:
            slide1_blocks_correct = 1
        if slides_contents[1]['titleCounter'] >= 1 or \
                (slides_contents[1]['textCounter'] >= 1 and not slides_contents[1]['titleCounter']):
            analyze_params[2] = True
        # End first slide #

        # Check second slide #
        slide2_font_size = 0
        slide2_blocks_correct = 0
        slide2_title = 0
        if max(font_sizes[2]) == 24 and min(font_sizes[2]) == 20:
            slide2_font_size = 1
        if slides_contents[2]['textCounter'] + slides_contents[2]['titleCounter'] == 3 and \
                slides_contents[2]['pictureCounter'] == 2:
            slide2_blocks_correct = 1
            slide2_title = 1
        if slides_contents[2]['textCounter'] + slides_contents[2]['titleCounter'] == 2 and \
                slides_contents[2]['pictureCounter'] == 2:
            slide2_blocks_correct = 1
        # End second slide #

        # Check third slide #
        slide3_font_size = 0
        slide3_blocks_correct = 0
        slide3_title = 0
        if max(font_sizes[3]) == 24 and min(font_sizes[3]) == 20:
            slide3_font_size = 1
        if slides_contents[3]['textCounter'] + slides_contents[3]['titleCounter'] == 3 and \
                slides_contents[3]['pictureCounter'] == 3:
            slide3_blocks_correct = 1
        if slides_contents[3]['textCounter'] + slides_contents[3]['titleCounter'] == 4 and \
                slides_contents[3]['pictureCounter'] == 3:
            slide3_blocks_correct = 1
            slide3_title = 1
        if slides_contents[3]['titleCounter'] == 1:
            slide3_title = 1
        # End third slide #

        if (slide1_blocks_correct + slide2_blocks_correct + slide3_blocks_correct) == 3:
            analyze_params[1] = True
        else:
            analyze_params[1] = False
        if (slide2_title + slide3_title) == 2:
            analyze_params[3] = True
        else:
            analyze_params[3] = False
        if (slide1_font_size + slide2_font_size + slide3_font_size) == 3:
            analyze_params[6] = True
        else:
            analyze_params[6] = False
        if not len(typefaces) > 1:
            analyze_params[5] = True
        else:
            analyze_params[5] = False
        return self._translate_results(analyze_params)


class CheckPresentationImages(CheckPresentationGetData):
    def generate_skeleton(self, flag=None):
        """
        Генерирует и удаляет показательные изображения коллизий.
        :param flag: Список путей файлов которые необходимо удалить
        :type flag: list of string
        :return: список путей если флаг None, True в случае успешного удаления иначе
        :rtype: list of string or boolean
        """
        if flag is None:
            result = []
            for slide in self.presentation.slides:
                index = int(self.presentation.slides.index(slide) + 1)
                im = Image.new(mode="RGB", size=self.prs_w_h())
                draw = ImageDraw.Draw(im)
                # previous_dimensions = []
                # previous_shape = None
                for shape in slide.shapes:
                    dim = self.dimensions_to_draw(shape)
                    if self.is_text(shape):
                        dim = self.dimensions_based_on_font(shape)
                    current_dimensions = [dim['x1'], dim['y1'], dim['width'], dim['height']]
                    image_color = "blue"
                    text_color = "yellow"
                    # if len(previous_dimensions) > 0 and previous_shape is not None:
                    # if self.check_collision(current_dimensions, previous_dimensions):
                    # print(self.check_collision(current_dimensions, previous_dimensions), shape.name,
                    # previous_shape.name)
                    if self.is_text(shape):
                        draw.rectangle([(dim['x1'], dim['y1']), (dim['x2'], dim['y2'])], fill=text_color, outline="red")
                    if self.is_image(shape):
                        draw.rectangle([(dim['x1'], dim['y1']), (dim['x2'], dim['y2'])], fill=image_color)
                    # previous_dimensions = current_dimensions
                    # previous_shape = shape
                im.save(f'{index}.jpg')
                result.append(f'{index}.jpg')
            return result
        else:
            for f in flag:
                os.mkdir(f)
            return True

    def save_images(self):
        """
        :return: Пути к сохранённым картинкам и их координаты Top Left на слайде
        :rtype: object of list of tuple (string, tuple of (int, int)
        """
        slide_counter, image_cords, image_paths = 0, [], []
        for slide in self.presentation.slides:
            slide_counter += 1
            picture_counter = 0
            for shape in slide.shapes:
                if self.is_image(shape):
                    picture_counter += 1
                    temp_image_paths = [
                        f"slide{slide_counter}_{picture_counter}.png",
                        f"slide{slide_counter}_{picture_counter}_original.png",
                    ]
                    image_size = (self.convert_emu_px(shape.width), self.convert_emu_px(shape.height))
                    original_image = open(temp_image_paths[1], 'wb')
                    original_image.write(base64.b64decode(base64.b64encode(shape.image.blob)))
                    original_image.close()
                    Image.open(temp_image_paths[1]).resize(image_size, Image.ANTIALIAS).save(temp_image_paths[0])
                    os.remove(temp_image_paths[1])
                    image_cords.append((self.convert_emu_px(shape.left), self.convert_emu_px(shape.top)))
                    image_paths.append(temp_image_paths[0])
        return list(zip(image_paths, image_cords))

    def generate_slide_images(self, bg="white", text_fill="black"):
        """
        Генерирует максимально точно возможные скриншоты слайдов не используя API PowerPoint. Используя убедитесь что
        у пользователя есть доступ к папке C:\\Windows\\Fonts\\, и в том что установлены стандартные шрифты.
        Изображения генерируются в папке с файлом.
        :param bg: Желаемый цвет заднего фона
        :type bg: string
        :param text_fill: Желаемый цвет текста
        :return: Список с путями к сохранённым фотографиям
        :rtype: list of string
        """
        result = []
        images = self.save_images()
        prs_size = self.prs_w_h()
        for slide in self.presentation.slides:
            idx = self.presentation.slides.index(slide) + 1
            slide_images = [(p[0], p[1]) for p in images if f'slide{idx}' in p[0]]  # p[0]: path ; p[1] width, height
            screen = Image.new("RGBA", prs_size, bg)
            for shape in slide.shapes:
                if self.is_text(shape):
                    dims = self.dimensions_to_draw(shape)
                    typeface = self.get_typefaces_by_shape(shape).split(' ')[0].lower()
                    font_size = int(max(self.font_sizes_by_shape(shape, True))) + 4
                    try:
                        font = ImageFont.truetype(font=f"C:/Windows/Fonts/{typeface}.ttf", size=font_size)
                    except OSError:
                        font = ImageFont.truetype(font=f"C:/Windows/Fonts/{DEFAULT_FONT.lower()}.ttf", size=font_size)
                    wrap = TextWrapper(shape.text, font, dims['width'])
                    text_image = Image.new("RGB", wrap.total_width_height(), bg)
                    draw = ImageDraw.Draw(text_image)
                    draw.text((0, 0), wrap.wrapped_text(), font=font, fill=text_fill)
                    screen.paste(text_image, (dims['x1'], dims['y1']))
                    # text_image.save(f'tmp/{self.random_string()}.jpg')
                    # process images
            for img in slide_images:
                screen.paste(Image.open(img[0]), img[1])
                os.remove(img[0])
            screen.save(f'slide{idx}.png')
            result.append(f'{os.getcwd()}\\slide{idx}.png')
        return result


class CheckPresentationPossibleWarnings(CheckPresentationGetData):
    def check_font_sizes(self):
        """
        Проверяет не установлено ли в каком либо текстовом обьекте значение по умолчанию. В таком случае значение
        нельзя получить в данной реализации программы: https://github.com/scanny/python-pptx/issues/337
        :return: False если все значения шрифта корректны, иначе список с некорректными классами.
        :rtype: list of class
        """
        result = []
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if self.is_text(shape):
                    if len(self.font_sizes_by_shape(shape, False)) == 0:
                        result.append(shape)
        if len(result) == 0:
            return False
        return result

    def check_slides(self):
        """
        Проверяет если в презентации больше трёх слайдов.
        :return: True если больше 3х слайдов иначе False
        :rtype: bool
        """
        if self.length() > 3:
            return True
        return False

    def check_crop_image(self):
        """
        Проверяет обрезалось ли изображение. Использует константное значение POSSIBLE_CROP_VALUE_ERROR == 0.05,
        что равно 5% обрезанного. Что укладывается в погрешность.
        :return: Обьект картинки в презенатции
        :rtype: list of class
        """
        result = []
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if self.is_image(shape):
                    if shape.crop_bottom > POSSIBLE_CROP_VALUE_ERROR or \
                            shape.crop_left > POSSIBLE_CROP_VALUE_ERROR or \
                            shape.crop_right > POSSIBLE_CROP_VALUE_ERROR or \
                            shape.crop_top > POSSIBLE_CROP_VALUE_ERROR:
                        result.append(shape)
        if len(result) == 0:
            return False
        return result

    def get_all(self):
        """
        Вернёт все предупреждения найденные в презентации
        :return: Строчки предупреждений
        :rtype: list of string
        """
        result = []
        if self.check_font_sizes():
            for shape in self.check_font_sizes():
                result.append(f'Предупреждение: в объекте презентации {shape.name} не удалось получить шрифт. При '
                              f'проверке использовано значение по умолчанию.')
        if self.check_slides():
            result.append('Предупреждение: количество слайдов больше трёх. Проверяются только первые три слайда. '
                          'Проверка может быть некорректной.')
        if self.check_crop_image():
            for shape in self.check_crop_image():
                text = f'Предупреждение: изображение {shape.name} было обрезано пользователем. Детали: \n'
                if shape.crop_bottom:
                    text += f"Снизу обрезано: {shape.crop_bottom * 100}%; \n"
                if shape.crop_top:
                    text += f"Снизу обрезано: {shape.crop_top * 100}%; \n"
                if shape.crop_left:
                    text += f"Снизу обрезано: {shape.crop_left * 100}%; \n"
                if shape.crop_right:
                    text += f"Снизу обрезано: {shape.crop_right * 100}%; \n"
                result.append(text)
        return result


class CheckPresentationTesting(CheckPresentationImages):
    def generate_slide_images(self):
        if self.path_to_presentation is not None:
            app = win32com.client.Dispatch("PowerPoint.Application")
            prs = app.Presentations.Open(self.path_to_presentation, WithWindow=False)
            directory = os.path.abspath(os.getcwd())
            counter = 1
            result = []
            for s in prs.Slides:
                path = f"{directory}\\slide_images\\slide_{counter}.jpg"
                s.Export(path, "JPG")
                result.append(path)
                counter += 1
            '''
            Этот фрагмент показывает как я пытался починить рамки shape'оф программно. Безуспешно пока что.
            for sld in prs.Slides:
                for shp in sld.Shapes:
                    if shp.HasTextFrame:
                        if not shp.TextFrame.AutoSize == MSOPPT.constants.ppAutoSizeShapeToFitText:
                            shp.TextFrame.WordWrap = MSO.constants.msoFalse
                            shp.TextFrame.AutoSize = MSOPPT.constants.ppAutoSizeShapeToFitText
                            shp.TextFrame.WordWrap = MSO.constants.msoTrue
            prs.Save()
            '''
            return result
        raise Exception("Не указан путь к презентации")
