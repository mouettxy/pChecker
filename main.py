import sys
import form
import pygame
import base64
import os
import shutil
import PIL
import re

from collections import Counter
from PIL import Image
from PyQt5.Qt import QMainWindow, QApplication, QFileInfo, QPixmap
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.enum.text import MSO_AUTO_SIZE


class PCheckerUtils:
    def __init__(self):
        super().__init__()
        self.mso_pic = MSO_SHAPE_TYPE.PICTURE
        self.placeholder_pic = PP_PLACEHOLDER_TYPE.PICTURE

    @staticmethod
    def emu_px(emu):
        return round(emu // 9525)

    def get_width_height(self, presentation):
        return self.emu_px(presentation.slide_width), self.emu_px(presentation.slide_height)

    def save_images(self, presentation):
        slide_counter, image_cords, image_paths = 0, [], []
        try:  # creates dir, if exist delete and recreate #
            os.mkdir('img')
        except OSError:
            shutil.rmtree('img', ignore_errors=True)
            os.mkdir('img')
        for slide in presentation.slides:
            slide_counter += 1
            picture_counter = 1
            for shape in slide.shapes:
                if shape.shape_type == self.mso_pic or (shape.is_placeholder and shape.placeholder_format.type ==
                                                                                 self.placeholder_pic):
                    if slide_counter > 3:
                        raise Exception('Количество слайдов более чем 3, остановка парсинга картинок')
                    pil_pic_path = f"img/img_slide{slide_counter}_{picture_counter}.png"
                    pic_path     = f"img/img_slide{slide_counter}_{picture_counter}_original.png"
                    picture_counter += 1
                    with open(pic_path, 'wb') as f:
                        f.write(base64.b64decode(base64.b64encode(shape.image.blob)))
                        f.close()
                    pic_size = (self.emu_px(shape.width), self.emu_px(shape.height))
                    Image.open(pic_path).resize(pic_size, Image.ANTIALIAS).save(pil_pic_path)
                    os.remove(pic_path)
                    if 1 < slide_counter < 4:
                        image_cords.append((self.emu_px(shape.left), self.emu_px(shape.top)))
                        image_paths.append(pil_pic_path)

                else:
                    pass
        return [[image_paths[i], image_cords[i]] for i in range(len(image_cords))]

    def get_cords_dim(self, presentation):
        slide_counter = 0
        cords_dim = {'2': [], '3': []}
        for slide in presentation.slides:
            slide_counter += 1
            for shape in slide.shapes:
                if 1 < slide_counter < 4:
                    if hasattr(shape, "text"):
                        left_top     = (self.emu_px(shape.left), self.emu_px(shape.top))
                        width_height = (self.emu_px(shape.width), self.emu_px(shape.height))
                        cords_dim[str(slide_counter)].append([left_top, width_height])
        return cords_dim

    @staticmethod
    def create_screenshots(window_size, image_path_cords, prs_cords_dim_text):
        try:
            os.mkdir('screens')
        except OSError:
            shutil.rmtree('screens')
            os.mkdir('screens')
        slides_pic = {'2': [path for path in image_path_cords if 'slide2' in path[0]],
                      '3': [path for path in image_path_cords if 'slide3' in path[0]]}
        slides_nums = (2, 3)
        for slide in slides_pic:
            pygame.init()
            screen = pygame.display.set_mode(window_size)
            for text_cords in prs_cords_dim_text[slide]:
                pygame.draw.rect(screen, (255, 255, 255), (text_cords[0][0], text_cords[0][1],
                                                           text_cords[1][0], text_cords[1][1]), 1)
            for path_cords in slides_pic[slide]:
                screen.blit(pygame.image.load(path_cords[0]), path_cords[1])
            pygame.display.update()
            pygame.image.save(screen, f"screens/screen{slide}.jpg")
            pygame.quit()
        return f"screens/screen{slides_nums[0]}.jpg", f"screens/screen{slides_nums[1]}.jpg"


class PChecker:
    def __init__(self, presentation):
        super().__init__()
        self.presentation  = presentation
        self.text_threshold = 2
        self.placeholder = PP_PLACEHOLDER_TYPE
        self.mso = MSO_SHAPE_TYPE
        self.warnings = []

    def generate_warnings(self, warning):
        self.warnings.append(warning)
        print(warning)

    def get_slides(self):
        return len(self.presentation.slides)

    def is_text(self, shape):
        if shape.has_text_frame:
            if self.is_title(shape):
                return True
            if shape.is_placeholder and shape.placeholder_format.type == self.placeholder.BODY:
                return True
            if len(shape.text) > self.text_threshold:
                return True
        return False

    def is_image(self, shape):
        if shape.shape_type == self.mso.PICTURE:
            return True
        if shape.is_placeholder and shape.placeholder_format.type == self.placeholder.PICTURE:
            return True
        return False

    def is_title(self, shape):
        if shape.is_placeholder and (
            shape.placeholder_format.type == self.placeholder.TITLE
                or shape.placeholder_format.type == self.placeholder.SUBTITLE
                or shape.placeholder_format.type == self.placeholder.VERTICAL_TITLE
                or shape.placeholder_format.type == self.placeholder.CENTER_TITLE):
            return True
        return False

    def _get_font_sizes_by_id(self, slide_id):
        font_sizes = []
        for shape in self.presentation.slides.get(slide_id).shapes:
            if self.is_text(shape):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        try:
                            if shape.text_frame.auto_size is None or \
                                    shape.text_frame.auto_size == MSO_AUTO_SIZE.NONE or \
                                    shape.text_frame.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT:
                                font_sizes.append(run.font.size.pt)
                            else:
                                font_sizes.append(run.font.size.pt *
                                                         shape.text_frame._bodyPr.normAutofit.fontScale / 100)
                        except AttributeError:
                            pass
        return font_sizes

    def _get_all_paragraph_runs(self):
        runs = []
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if self.is_text(shape):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            runs.append(run)
        return runs

    def get_font_sizes(self):
        font_sizes = {
        }
        for slide in self.presentation.slides:
            index = int(self.presentation.slides.index(slide) + 1)
            if index > 3:
                return font_sizes
            font_sizes.update({index: self._get_font_sizes_by_id(slide.slide_id)})
        return font_sizes

    def get_typefaces(self):
        typefaces = set()
        for run in self._get_all_paragraph_runs():
            try:
                typefaces.add(run.font.name)
            except AttributeError:
                pass
        return typefaces

    def get_text(self):
        text = []
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if self.is_text(shape):
                    text.append(shape.text)
        return text

    @staticmethod
    def string_optimize(string):
        """
        Method reproduce given string and deletes all words that length <= 3 and all punctuation symbols
        """
        delete_junk_symbols = re.compile('[^a-zA-Zа-яА-ЯёЁ ]')
        delete_junk_words = re.compile('\\b\\w{0,3}\\b')
        return delete_junk_words.sub("", delete_junk_symbols.sub("", string.lower()))

    def analyze_text(self):
        text_analyzed = []
        for text in self.get_text():
            text_analyzed.extend(self.string_optimize(text).split())
        most_common = Counter(text_analyzed).most_common(5)
        return most_common

    '''Lets check it out *On background starts song AC/DC - Highway to Hell*'''

    def get_slides_contents(self):
        slides = {
            1: {
                'textCounter': 0,
                'pictureCounter': 0,
                'titleCounter': 0,
            },
            2: {
                'textCounter': 0,
                'pictureCounter': 0,
                'titleCounter': 0,
            },
            3: {
                'textCounter': 0,
                'pictureCounter': 0,
                'titleCounter': 0,
            }
        }
        for slide in self.presentation.slides:
            index = int(self.presentation.slides.index(slide) + 1)
            if index > 3:
                return slides
            for shape in slide.shapes:
                if self.is_text(shape):
                    if self.is_title(shape):
                        slides[index]['titleCounter'] += 1
                    else:
                        slides[index]['textCounter'] += 1
                if self.is_image(shape):
                    slides[index]['pictureCounter'] += 1
        return slides

    def analyze_results(self,
                        txt_img_collisions_btn=False,
                        distorted_images_btn=False,
                        all_collisions_btn=False,
                        content_compliance=False
                        ):
        analyze_params = {
            'slides_count':               self.get_slides(),
            'text_blocks_exist':          None,
            'title_on_cover_page':        None,
            'title_on_other_slides':      None,
            'content_compliance':         content_compliance,
            'single_typeface':            None,
            'right_font_size':            None,
            'text_not_overlaps_images':   txt_img_collisions_btn,
            'images_not_distorted':       distorted_images_btn,
            'images_not_overlaps_shapes': all_collisions_btn,
        }
        slides_contents = self.get_slides_contents()
        font_sizes      = self.get_font_sizes()
        typefaces       = self.get_typefaces()
        # Check first slide #
        slide1_font_size = 0
        slide1_blocks_correct = 0
        if max(font_sizes[1]) == 40 and min(font_sizes[1]) == 24:
            slide1_font_size = 1
        if slides_contents[1]['titleCounter'] + slides_contents[1]['textCounter'] == 2:
            slide1_blocks_correct = 1
        if slides_contents[1]['titleCounter'] >= 1 or \
                (slides_contents[1]['textCounter'] >= 1 and not slides_contents[1]['titleCounter']):
            analyze_params['title_on_cover_page'] = True
        # End first slide #

        # Check second slide #
        slide2_font_size = 0
        slide2_blocks_correct = 0
        slide2_title = 0
        if max(font_sizes[2]) == 24 and min(font_sizes[2]) == 20:
            slide2_font_size = 1
        if slides_contents[2]['textCounter'] + slides_contents[2]['titleCounter'] == 3 and \
           slides_contents[2]['pictureCounter'] == 2:
            slide2_blocks_correct = 1
            slide2_title = 1
        if slides_contents[2]['textCounter'] + slides_contents[2]['titleCounter'] == 2 and \
           slides_contents[2]['pictureCounter'] == 2:
            slide2_blocks_correct = 1
        # End second slide #

        # Check third slide #
        slide3_font_size = 0
        slide3_blocks_correct = 0
        slide3_title = 0
        if max(font_sizes[3]) == 24 and min(font_sizes[3]) == 20:
            slide3_font_size = 1
        if slides_contents[3]['textCounter'] + slides_contents[3]['titleCounter'] == 3 and \
           slides_contents[3]['pictureCounter'] == 3:
            slide3_blocks_correct = 1
        if slides_contents[3]['textCounter'] + slides_contents[3]['titleCounter'] == 4 and \
           slides_contents[3]['pictureCounter'] == 3:
            slide3_blocks_correct = 1
            slide3_title = 1
        if slides_contents[3]['titleCounter'] == 1:
            slide3_title = 1
        # End third slide #

        if (slide1_blocks_correct + slide2_blocks_correct + slide3_blocks_correct) == 3:
            analyze_params['text_blocks_exist'] = True
        else:
            analyze_params['text_blocks_exist'] = False
        if (slide2_title + slide3_title) == 2:
            analyze_params['title_on_other_slides'] = True
        else:
            analyze_params['title_on_other_slides'] = False
        if (slide1_font_size + slide2_font_size + slide3_font_size) == 3:
            analyze_params['right_font_size'] = True
        else:
            analyze_params['right_font_size'] = False
        if not len(typefaces) > 1:
            analyze_params['single_typeface'] = True
        else:
            analyze_params['single_typeface'] = False
        return analyze_params

    @staticmethod
    def translate_results(results):
        results['Количество слайдов']                   = results.pop('slides_count')
        results['Блоки текста и изображений размещены'] = results.pop('text_blocks_exist')
        results['Название на титульном']                = results.pop('title_on_cover_page')
        results['Название на 2м и 3м слайде']           = results.pop('title_on_other_slides')
        results['Соответствие теме']                    = results.pop('content_compliance')
        results['Единый шрифт']                         = results.pop('single_typeface')
        results['Правильный размер шрифта']             = results.pop('right_font_size')
        results['Текст не перекрывает изображения']     = results.pop('text_not_overlaps_images')
        results['Изображения не искажены']              = results.pop('images_not_distorted')
        results['Изображения не перекрывают элементы']  = results.pop('images_not_overlaps_shapes')
        return results

    def unloading_xlsx(self):
        pass

    def unloading_txt(self):
        pass


class PCheckerWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.Utils = PCheckerUtils()
        self.ui = form.Ui_MainWindow()
        self.ui.setupUi(self)
        self.ui.get_answer.clicked.connect(self.checkAndReturnResult)
        self.setAcceptDrops(True)

    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls():
            event.accept()
        else:
            event.ignore()

    def dropEvent(self, event):  # Ловим ивент дропа файла в окно
        self.file = event.mimeData().urls()[0].toLocalFile()
        file_extension = QFileInfo(self.file).suffix()
        try:
            if file_extension == 'jpg' or file_extension == 'jpeg' or file_extension == 'png':
                self.ui.image_holder.setPixmap(QPixmap(self.file))
                self.ui.statusbar.append(f'Изображение по пути {self.file} установлено')
            elif file_extension == 'pptx':
                Check = PChecker(Presentation(self.file))
                slide_size         = self.Utils.get_width_height(Presentation(self.file))
                images_path_cords  = self.Utils.save_images(Presentation(self.file))
                prs_cords_dim_text = self.Utils.get_cords_dim(Presentation(self.file))
                screens            = self.Utils.create_screenshots(slide_size, images_path_cords, prs_cords_dim_text)
                content_compliance = Check.analyze_text()
                content_compliance = [words for words in content_compliance]
                self.ui.slide2_image_label.setPixmap(QPixmap(screens[0]))
                self.ui.slide3_image_label.setPixmap(QPixmap(screens[1]))
                self.ui.statusbar.append(f'Разбор pptx файла по пути {self.file}')
                self.ui.statusbar.append(f'Наиболее часто встречающиеся слова в презентации:')
                for tuple_content in content_compliance:
                    self.ui.statusbar.append(f'{tuple_content[0]} => {tuple_content[1]}')
            else:
                self.ui.statusbar.append('Не поддерживаемое расширение файла.')
        except Exception as e:
            self.ui.statusbar.append(f'Ошибка {e} при загрузке файла.')

    def checkAndReturnResult(self):
        Check = PChecker(Presentation(self.file))
        img_coll_btn = self.ui.txt_img_collisions_btn_yes.isChecked()
        dist_images = self.ui.distorted_images_btn_yes.isChecked()
        all_coll = self.ui.all_collisions_btn_yes.isChecked()
        cont_compliance = self.ui.content_compliance_btn_yes.isChecked()
        unloading_txt = self.ui.unloading_btn1.isChecked()
        unloading_xlsx = self.ui.unloading_btn2.isChecked()
        unloading_statusbar = self.ui.unloading_btn3.isChecked()
        results = Check.analyze_results(img_coll_btn, dist_images, all_coll, cont_compliance)
        translated_result = Check.translate_results(results)
        if unloading_statusbar:
            for result in translated_result:
                self.ui.statusbar.append(f'{result} === {translated_result[result]}')
        elif unloading_txt:
            Check.unloading_txt()
        else:
            Check.unloading_xlsx()


if __name__ == "__main__":
    app = QApplication(sys.argv)
    ex = PCheckerWindow()
    ex.show()
    sys.exit(app.exec_())
