import os
import base64
import pygame
import re
from shutil import rmtree
from PIL import Image
from collections import Counter
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.enum.text import MSO_AUTO_SIZE


class PresentationCustomUtils:
    def __init__(self):
        super().__init__()

    @staticmethod
    def emu_px(emu):
        return round(emu // 9525)

    @staticmethod
    def create_screenshots(window_size, image_path_cords, prs_cords_dim_text):
        try:
            os.mkdir('screens')
        except OSError:
            rmtree('screens', ignore_errors=True)
            os.mkdir('screens')
        slides_pic = {'2': [path for path in image_path_cords if 'slide2' in path[0]],
                      '3': [path for path in image_path_cords if 'slide3' in path[0]]}
        slides_nums = (2, 3)
        for slide in slides_pic:
            pygame.init()
            screen = pygame.display.set_mode(window_size)
            for text_cords in prs_cords_dim_text[slide]:
                pygame.draw.rect(screen, (255, 255, 255), (text_cords[0][0], text_cords[0][1],
                                                           text_cords[1][0], text_cords[1][1]), 1)
            for path_cords in slides_pic[slide]:
                screen.blit(pygame.image.load(path_cords[0]), path_cords[1])
            pygame.display.update()
            pygame.image.save(screen, f"screens/screen{slide}.jpg")
            pygame.quit()
        return f"screens/screen{slides_nums[0]}.jpg", f"screens/screen{slides_nums[1]}.jpg"

    @staticmethod
    def string_optimize(string):
        """
        Method reproduce given string and deletes all words that length <= 3 and all punctuation symbols
        """
        delete_junk_symbols = re.compile('[^a-zA-Zа-яА-ЯёЁ ]')
        delete_junk_words = re.compile('\\b\\w{0,3}\\b')
        return delete_junk_words.sub("", delete_junk_symbols.sub("", string.lower()))

    @staticmethod
    def translate_results(results):
        results['Количество слайдов'] = results.pop('slides_count')
        results['Блоки текста и изображений размещены'] = results.pop('text_blocks_exist')
        results['Название на титульном'] = results.pop('title_on_cover_page')
        results['Название на 2м и 3м слайде'] = results.pop('title_on_other_slides')
        results['Соответствие теме'] = results.pop('content_compliance')
        results['Единый шрифт'] = results.pop('single_typeface')
        results['Правильный размер шрифта'] = results.pop('right_font_size')
        results['Текст не перекрывает изображения'] = results.pop('text_not_overlaps_images')
        results['Изображения не искажены'] = results.pop('images_not_distorted')
        results['Изображения не перекрывают элементы'] = results.pop('images_not_overlaps_shapes')
        return results

    def get_width_height(self, presentation):
        return self.emu_px(presentation.slide_width), self.emu_px(presentation.slide_height)

    def save_images(self, presentation):
        slide_counter, image_cords, image_paths = 0, [], []
        try:  # creates dir, if exist delete and recreate #
            os.mkdir('img')
        except OSError:
            rmtree('img', ignore_errors=True)
            os.mkdir('img')
        for slide in presentation.slides:
            slide_counter += 1
            picture_counter = 1
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (shape.is_placeholder and
                                                                  shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.PICTURE):
                    if slide_counter > 3:
                        raise Exception('Количество слайдов более чем 3, остановка парсинга картинок')
                    pil_pic_path = f"img/img_slide{slide_counter}_{picture_counter}.png"
                    pic_path     = f"img/img_slide{slide_counter}_{picture_counter}_original.png"
                    picture_counter += 1
                    with open(pic_path, 'wb') as f:
                        f.write(base64.b64decode(base64.b64encode(shape.image.blob)))
                        f.close()
                    pic_size = (self.emu_px(shape.width), self.emu_px(shape.height))
                    Image.open(pic_path).resize(pic_size, Image.ANTIALIAS).save(pil_pic_path)
                    os.remove(pic_path)
                    if 1 < slide_counter < 4:
                        image_cords.append((self.emu_px(shape.left), self.emu_px(shape.top)))
                        image_paths.append(pil_pic_path)
                else:
                    pass
        return [[image_paths[i], image_cords[i]] for i in range(len(image_cords))]

    def get_cords_dim(self, presentation):
        slide_counter = 0
        cords_dim = {'2': [], '3': []}
        for slide in presentation.slides:
            slide_counter += 1
            for shape in slide.shapes:
                if 1 < slide_counter < 4:
                    if hasattr(shape, "text"):
                        left_top     = (self.emu_px(shape.left), self.emu_px(shape.top))
                        width_height = (self.emu_px(shape.width), self.emu_px(shape.height))
                        cords_dim[str(slide_counter)].append([left_top, width_height])
        return cords_dim


class PresentationUtils(PresentationCustomUtils):
    def __init__(self, presentation):
        super().__init__()
        self.presentation  = presentation
        self.text_threshold = 2

    def get_slides_len(self):
        return len(self.presentation.slides)

    def is_text(self, shape):
        if shape.has_text_frame:
            if self.is_title(shape):
                return True
            if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.BODY:
                return True
            if len(shape.text) > self.text_threshold:
                return True
        return False

    def is_image(self, shape):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.PICTURE:
            return True
        return False

    def is_title(self, shape):
        if shape.is_placeholder and (
            shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.TITLE
                or shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.SUBTITLE
                or shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.VERTICAL_TITLE
                or shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.CENTER_TITLE):
            return True
        return False

    def _get_font_sizes_by_id(self, slide_id):
        font_sizes = []
        for shape in self.presentation.slides.get(slide_id).shapes:
            if self.is_text(shape):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        try:
                            if shape.text_frame.auto_size is None or \
                                    shape.text_frame.auto_size == MSO_AUTO_SIZE.NONE or \
                                    shape.text_frame.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT:
                                font_sizes.append(run.font.size.pt)
                            else:
                                font_sizes.append(run.font.size.pt *
                                                         shape.text_frame._bodyPr.normAutofit.fontScale / 100)
                        except AttributeError:
                            pass
        return font_sizes

    def _get_all_paragraph_runs(self):
        runs = []
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if self.is_text(shape):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            runs.append(run)
        return runs

    def get_font_sizes(self):
        font_sizes = {}
        for slide in self.presentation.slides:
            index = int(self.presentation.slides.index(slide) + 1)
            if index > 3:
                return font_sizes
            font_sizes.update({index: self._get_font_sizes_by_id(slide.slide_id)})
        return font_sizes

    def get_typefaces(self):
        typefaces = set()
        for run in self._get_all_paragraph_runs():
            try:
                typefaces.add(run.font.name)
            except AttributeError:
                pass
        return typefaces

    def get_text(self):
        text = []
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if self.is_text(shape):
                    text.append(shape.text)
        return text

    def analyze_text(self):
        text_analyzed = []
        for text in self.get_text():
            text_analyzed.extend(self.string_optimize(text).split())
        most_common = Counter(text_analyzed).most_common(5)
        return most_common